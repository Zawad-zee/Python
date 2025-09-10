import cv2
import numpy as np
import os
from datetime import datetime
import glob

# Try to import PowerPoint library
try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PowerPoint generation will be disabled.")
    print("Install it using: pip install python-pptx")


class VideoFrameDetector:
    def __init__(self, video_path, output_dir="snapshots", threshold=30.0, min_interval=1.0, debug=False,
                 buffer_time=2.0):
        """
        Initialize the video frame change detector.

        Args:
            video_path (str): Path to the input video file
            output_dir (str): Directory to save snapshots
            threshold (float): Sensitivity threshold for detecting changes (0-100)
            min_interval (float): Minimum seconds between snapshots
            debug (bool): Print detailed information about frame changes
            buffer_time (float): Wait time after detecting change before taking snapshot
        """
        self.video_path = video_path
        self.output_dir = output_dir
        self.threshold = threshold
        self.min_interval = min_interval
        self.debug = debug
        self.buffer_time = buffer_time
        self.last_snapshot_time = 0
        self.pending_snapshots = []  # Store changes that need buffer time
        self.snapshot_files = []  # Keep track of saved snapshots

        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)

    def calculate_frame_difference(self, frame1, frame2):
        """
        Calculate the percentage difference between two frames with text sensitivity.

        Args:
            frame1, frame2: OpenCV frames to compare

        Returns:
            float: Percentage difference between frames
        """
        # Convert frames to grayscale for comparison
        gray1 = cv2.cvtColor(frame1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(frame2, cv2.COLOR_BGR2GRAY)

        # Calculate absolute difference
        diff = cv2.absdiff(gray1, gray2)

        # Method 1: Standard pixel counting (for major changes)
        total_pixels = diff.shape[0] * diff.shape[1]
        changed_pixels = np.count_nonzero(diff > 25)
        standard_percentage = (changed_pixels / total_pixels) * 100

        # Method 2: Text-sensitive detection (for small changes)
        # Count pixels with even small changes (better for text)
        text_sensitive_pixels = np.count_nonzero(diff > 10)
        text_percentage = (text_sensitive_pixels / total_pixels) * 100

        # Method 3: Edge detection for text changes
        edges1 = cv2.Canny(gray1, 50, 150)
        edges2 = cv2.Canny(gray2, 50, 150)
        edge_diff = cv2.absdiff(edges1, edges2)
        edge_changes = np.count_nonzero(edge_diff > 0)
        edge_percentage = (edge_changes / total_pixels) * 100

        # Use the maximum of all methods to catch both major and minor changes
        return max(standard_percentage, text_percentage * 0.5, edge_percentage * 0.3)

    def save_snapshot(self, frame, frame_number, timestamp):
        """
        Save a frame as a snapshot image.

        Args:
            frame: OpenCV frame to save
            frame_number (int): Frame number in the video
            timestamp (float): Video timestamp in seconds
        """
        # Create filename with timestamp
        time_str = datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]  # Include milliseconds
        filename = f"snapshot_{time_str}_frame{frame_number}_t{timestamp:.2f}s.jpg"
        filepath = os.path.join(self.output_dir, filename)

        # Save the frame
        cv2.imwrite(filepath, frame)
        self.snapshot_files.append(filepath)  # Track the saved file
        print(f"Snapshot saved: {filename}")

    def process_video(self):
        """
        Process the video and capture snapshots when frames change significantly.
        """
        # Open video file
        cap = cv2.VideoCapture(self.video_path)

        if not cap.isOpened():
            print(f"Error: Could not open video file {self.video_path}")
            return

        # Get video properties
        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps

        print(f"Processing video: {self.video_path}")
        print(f"FPS: {fps:.2f}, Duration: {duration:.2f}s, Total frames: {total_frames}")
        print(f"Change threshold: {self.threshold}%")
        print(f"Minimum interval: {self.min_interval}s, Buffer time: {self.buffer_time}s")
        print("-" * 50)

        # Read first frame
        ret, prev_frame = cap.read()
        if not ret:
            print("Error: Could not read first frame")
            return

        frame_count = 0
        snapshots_taken = 0

        # Save first frame as initial snapshot
        self.save_snapshot(prev_frame, frame_count, 0)
        snapshots_taken += 1

        while True:
            ret, current_frame = cap.read()
            if not ret:
                break

            frame_count += 1
            current_time = frame_count / fps

            # Calculate difference between current and previous frame
            diff_percentage = self.calculate_frame_difference(prev_frame, current_frame)

            # Debug: Print all frame differences if debug mode is on
            if self.debug and frame_count % 30 == 0:  # Print every 30th frame to avoid spam
                print(f"Frame {frame_count}: {diff_percentage:.2f}% change (threshold: {self.threshold}%)")

            # Check if change exceeds threshold and minimum interval has passed
            time_since_last = current_time - self.last_snapshot_time

            if (diff_percentage > self.threshold and
                    time_since_last >= self.min_interval):
                # Add to pending snapshots with buffer time
                self.pending_snapshots.append({
                    'frame': current_frame.copy(),
                    'frame_number': frame_count,
                    'detection_time': current_time,
                    'snapshot_time': current_time + self.buffer_time,  # Take snapshot after buffer
                    'diff_percentage': diff_percentage
                })

                print(
                    f"📋 CHANGE DETECTED: Frame {frame_count}, {diff_percentage:.2f}% change - will snapshot in {self.buffer_time}s")

            # Check if any pending snapshots are ready (buffer time has passed)
            ready_snapshots = []
            for i, pending in enumerate(self.pending_snapshots):
                if current_time >= pending['snapshot_time']:
                    ready_snapshots.append(i)

            # Take snapshots for ready ones (in reverse order to maintain list indices)
            for i in reversed(ready_snapshots):
                pending = self.pending_snapshots[i]

                # Use current frame (after buffer time) instead of the detection frame
                self.save_snapshot(current_frame, frame_count, current_time)
                self.last_snapshot_time = current_time
                snapshots_taken += 1

                buffer_waited = current_time - pending['detection_time']
                print(
                    f"✓ BUFFERED SNAPSHOT {snapshots_taken}: Frame {frame_count} (waited {buffer_waited:.1f}s for animations)")

                # Remove from pending list
                self.pending_snapshots.pop(i)

            # Update previous frame
            prev_frame = current_frame.copy()

            # Show progress every 100 frames
            if frame_count % 100 == 0:
                progress = (frame_count / total_frames) * 100
                print(f"Progress: {progress:.1f}% ({frame_count}/{total_frames} frames)")

        cap.release()

        # Handle any remaining pending snapshots at the end
        if self.pending_snapshots:
            print(f"\nProcessing {len(self.pending_snapshots)} final pending snapshots...")
            for pending in self.pending_snapshots:
                self.save_snapshot(pending['frame'], pending['frame_number'], pending['detection_time'])
                snapshots_taken += 1
                print(f"✓ FINAL SNAPSHOT: Frame {pending['frame_number']} (end of video)")

        print("-" * 50)
        print(f"Processing complete!")
        print(f"Total snapshots taken: {snapshots_taken}")
        print(f"Snapshots saved in: {self.output_dir}")

        return snapshots_taken

    def create_powerpoint(self, pptx_filename=None):
        """
        Create a PowerPoint presentation from the saved snapshots.

        Args:
            pptx_filename (str): Optional custom filename for the PowerPoint file
        """
        if not PPTX_AVAILABLE:
            print("❌ PowerPoint creation failed: python-pptx library not installed")
            print("Install it using: pip install python-pptx")
            return False

        if not self.snapshot_files:
            print("❌ No snapshots found to create PowerPoint presentation")
            return False

        # Sort snapshot files by creation time (should already be in order)
        self.snapshot_files.sort()

        # Create PowerPoint presentation
        prs = Presentation()

        # Remove default slide layout
        if len(prs.slides) > 0:
            for slide in list(prs.slides):
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]

        print(f"📊 Creating PowerPoint with {len(self.snapshot_files)} slides...")

        for i, image_path in enumerate(self.snapshot_files, 1):
            if not os.path.exists(image_path):
                print(f"⚠️  Warning: Image not found: {image_path}")
                continue

            try:
                # Add a blank slide
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)

                # Calculate image dimensions to fit slide while maintaining aspect ratio
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # Add image to slide
                left = Inches(0.5)
                top = Inches(0.5)
                width = slide_width - Inches(1)  # Leave 0.5" margin on each side

                # Add the image
                pic = slide.shapes.add_picture(image_path, left, top, width=width)

                print(f"✓ Added slide {i}")

            except Exception as e:
                print(f"❌ Error adding slide {i}: {str(e)}")
                continue

        # Save PowerPoint file
        if not pptx_filename:
            video_name = os.path.splitext(os.path.basename(self.video_path))[0]
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_filename = f"{video_name}_snapshots_{timestamp}.pptx"

        pptx_path = os.path.join(self.output_dir, pptx_filename)

        try:
            prs.save(pptx_path)
            print(f"✅ PowerPoint presentation saved: {pptx_path}")
            return True
        except Exception as e:
            print(f"❌ Error saving PowerPoint: {str(e)}")
            return False


def main():
    # Auto-detect video files in current directory
    video_extensions = ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv']
    video_files = []

    for file in os.listdir('.'):
        if any(file.lower().endswith(ext) for ext in video_extensions):
            video_files.append(file)

    if not video_files:
        print("No video files found in current directory!")
        print("Please add a video file (.mp4, .avi, .mov, etc.) to this folder.")
        input("Press Enter to exit...")
        return

    if len(video_files) == 1:
        video_file = video_files[0]
        print(f"Found video file: {video_file}")
    else:
        print("Multiple video files found:")
        for i, file in enumerate(video_files, 1):
            print(f"{i}. {file}")

        while True:
            try:
                choice = int(input("Enter the number of the video to process: ")) - 1
                if 0 <= choice < len(video_files):
                    video_file = video_files[choice]
                    break
                else:
                    print("Invalid choice. Please try again.")
            except ValueError:
                print("Please enter a valid number.")

    # Ask for settings (with good defaults)
    print("\nSettings (press Enter for default values):")

    threshold = input("Sensitivity threshold (default 5.0, lower = more sensitive): ").strip()
    threshold = float(threshold) if threshold else 5.0

    buffer_time = input("Animation wait time in seconds (default 2.0): ").strip()
    buffer_time = float(buffer_time) if buffer_time else 2.0

    # Ask about PowerPoint creation
    create_pptx = True
    if PPTX_AVAILABLE:
        pptx_choice = input("Create PowerPoint presentation? (Y/n): ").strip().lower()
        create_pptx = pptx_choice != 'n'
    else:
        create_pptx = False
        print("PowerPoint creation disabled (python-pptx not installed)")

    print(f"\nProcessing {video_file} with threshold={threshold}, buffer_time={buffer_time}")
    print("-" * 50)

    # Create detector with user settings
    detector = VideoFrameDetector(
        video_path=video_file,
        output_dir="video_snapshots",
        threshold=threshold,
        min_interval=0.1,
        buffer_time=buffer_time,
        debug=True
    )

    # Process the video
    snapshots_taken = detector.process_video()

    # Create PowerPoint if requested and possible
    if create_pptx and snapshots_taken > 0:
        print("\n" + "=" * 50)
        print("Creating PowerPoint presentation...")
        detector.create_powerpoint()

    # Final summary
    print(f"\n" + "=" * 50)
    print(f"✅ PROCESSING COMPLETE!")
    print(f"📸 Total snapshots: {snapshots_taken}")
    print(f"📁 Snapshots saved in: {detector.output_dir}")

    if create_pptx and PPTX_AVAILABLE and snapshots_taken > 0:
        print("📊 PowerPoint presentation created!")

    input("\nPress Enter to exit...")


if __name__ == "__main__":
    # Check if OpenCV is installed
    try:
        import cv2

        print(f"OpenCV version: {cv2.__version__}")
        main()
    except ImportError:
        print("Error: OpenCV not installed.")
        print("Install it using: pip install opencv-python")